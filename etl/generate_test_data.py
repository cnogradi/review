from pptx import Presentation
from pptx.util import Inches
import os

def create_test_pptx(filename):
    prs = Presentation()
    
    # Slide 1
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created by Automated Test"

    # Slide 2
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Content Slide"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Bullet point 1"
    p = tf.add_paragraph()
    p.text = "Bullet point 2"
    p.level = 1

    prs.save(filename)
    print(f"Created {filename}")

if __name__ == "__main__":
    os.makedirs("test_docs", exist_ok=True)
    create_test_pptx("test_docs/test_presentation.pptx")
