import os
import json
import argparse
import shutil
from pathlib import Path
from typing import List, Dict, Any
import boto3
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from unstructured.partition.auto import partition
from PIL import Image
import io

def setup_args():
    parser = argparse.ArgumentParser(description='Process training documents and upload to S3')
    parser.add_argument('--source-dir', required=True, help='Directory containing source documents')
    parser.add_argument('--bucket', required=True, help='S3 bucket name')
    parser.add_argument('--s3-prefix', default='', help='Prefix for S3 keys')
    parser.add_argument('--dry-run', action='store_true', help='Skip S3 upload and save locally')
    parser.add_argument('--output-dir', default='./output', help='Local output directory for dry run')
    return parser.parse_args()

def extract_concepts(text: str) -> List[str]:
    # Simple heuristic: extract capitalized phrases or keywords
    # In a real scenario, this would use an LLM or NLP library
    words = text.split()
    concepts = [w.strip('.,()') for w in words if len(w) > 5 and w[0].isupper()]
    return list(set(concepts))[:5]  # Return top 5 unique concepts

def convert_pptx_to_pdf(pptx_path: Path, output_dir: Path) -> Path:
    """Converts PPTX to PDF using LibreOffice. Returns path to generated PDF."""
    import subprocess
    
    # Check for libreoffice or soffice
    soffice = shutil.which('soffice') or shutil.which('libreoffice')
    if not soffice:
        # On Windows, it might be in standard paths but not in PATH
        if os.name == 'nt':
            possible_paths = [
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
            ]
            for p in possible_paths:
                if os.path.exists(p):
                    soffice = p
                    break
    
    if not soffice:
        raise FileNotFoundError("LibreOffice (soffice) not found. Please install LibreOffice.")

    # Run conversion
    # --headless --convert-to pdf --outdir <dir> <file>
    cmd = [
        soffice,
        '--headless',
        '--convert-to', 'pdf',
        '--outdir', str(output_dir),
        str(pptx_path)
    ]
    
    subprocess.run(cmd, check=True, capture_output=True)
    
    pdf_name = pptx_path.with_suffix('.pdf').name
    return output_dir / pdf_name

def process_pptx(file_path: Path, output_dir: Path) -> Dict[str, Any]:
    prs = Presentation(file_path)
    slides_data = []
    images_dir = output_dir / 'images'
    images_dir.mkdir(parents=True, exist_ok=True)

    full_text = []
    
    # Try to generate slide previews using LibreOffice -> PDF -> Images
    slide_images = {}
    try:
        from pdf2image import convert_from_path
        
        # Create a temp dir for the PDF
        temp_pdf_dir = output_dir / 'temp_pdf'
        temp_pdf_dir.mkdir(exist_ok=True)
        
        try:
            print(f"Converting {file_path.name} to PDF for previews...")
            pdf_path = convert_pptx_to_pdf(file_path, temp_pdf_dir)
            
            print("Converting PDF to images...")
            # Convert PDF to images
            images = convert_from_path(str(pdf_path))
            
            for i, image in enumerate(images):
                image_filename = f"slide_{i+1}.jpg"
                image_path = images_dir / image_filename
                image.save(image_path, 'JPEG')
                slide_images[i] = f"images/{image_filename}"
                
        except Exception as e:
            print(f"Warning: Could not generate slide previews via LibreOffice: {e}")
            print("Ensure LibreOffice and poppler-utils are installed.")
        finally:
            if temp_pdf_dir.exists():
                shutil.rmtree(temp_pdf_dir)
                
    except ImportError:
        print("Warning: pdf2image not installed. Skipping slide preview generation.")

    for i, slide in enumerate(prs.slides):
        slide_text = []
        slide_title = f"Slide {i+1}"
        
        # Extract text and title using python-pptx (reliable for text)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_text.append(text)
                    if shape == slide.shapes.title:
                        slide_title = text

        # Get the generated image path or fallback
        slide_image_rel_path = slide_images.get(i)
        
        if not slide_image_rel_path:
            # Fallback: Try to extract embedded image if it's the main content
            image_filename = f"slide_{i+1}_embedded.jpg"
            image_path = images_dir / image_filename
            
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        with open(image_path, 'wb') as f:
                            f.write(image.blob)
                        slide_image_rel_path = f"images/{image_filename}"
                        break
                    except Exception:
                        pass

        content = "\n".join(slide_text)
        full_text.append(f"## {slide_title}\n\n{content}")
        
        slides_data.append({
            "id": f"slide-{i+1}",
            "title": slide_title,
            "level": 1,
            "concepts": extract_concepts(content),
            "preview_image": slide_image_rel_path
        })

    return {
        "outline": slides_data,
        "content": "\n\n".join(full_text)
    }

def process_unstructured(file_path: Path) -> Dict[str, Any]:
    elements = partition(filename=str(file_path))
    text_content = "\n\n".join([str(e) for e in elements])
    
    # Simple outline generation based on chunks
    # Unstructured elements might have types like 'Title'.
    outline = []
    current_section = None
    
    for i, element in enumerate(elements):
        if "Title" in str(type(element)):
            title = str(element)
            current_section = {
                "id": f"sec-{i}",
                "title": title,
                "level": 1,
                "concepts": []
            }
            outline.append(current_section)
        elif current_section:
            # Add concepts from text under this section
            current_section["concepts"].extend(extract_concepts(str(element)))
            # limit concepts
            current_section["concepts"] = list(set(current_section["concepts"]))[:5]

    if not outline:
        # Fallback if no titles found
        outline.append({
            "id": "doc-1",
            "title": "Main Content",
            "level": 1,
            "concepts": extract_concepts(text_content)
        })

    return {
        "outline": outline,
        "content": text_content
    }

def main():
    args = setup_args()
    source_dir = Path(args.source_dir)
    
    if args.dry_run:
        output_base = Path(args.output_dir)
        if output_base.exists():
            shutil.rmtree(output_base)
        output_base.mkdir(parents=True)
    else:
        s3 = boto3.client('s3')

    materials = []

    for file_path in source_dir.glob('*'):
        if file_path.name.startswith('.') or file_path.is_dir():
            continue
            
        print(f"Processing {file_path.name}...")
        
        module_code = file_path.stem.lower().replace(' ', '_')
        material_id = f"mat-{module_code}"
        
        # Prepare local output structure for this file
        if args.dry_run:
            module_output_dir = output_base / module_code
        else:
            module_output_dir = Path(f"/tmp/{module_code}") # Temp dir for S3 upload
            if module_output_dir.exists():
                shutil.rmtree(module_output_dir)
        
        module_output_dir.mkdir(parents=True, exist_ok=True)
        
        try:
            result = {}
            if file_path.suffix.lower() == '.pptx':
                result = process_pptx(file_path, module_output_dir)
            elif file_path.suffix.lower() in ['.pdf', '.docx', '.doc']:
                result = process_unstructured(file_path)
            else:
                print(f"Skipping unsupported file: {file_path.name}")
                continue

            material_data = {
                "id": material_id,
                "title": file_path.stem,
                "source": file_path.name,
                "outline": result["outline"],
                "content": result["content"]
            }
            
            # Save data.json
            with open(module_output_dir / 'data.json', 'w') as f:
                json.dump(material_data, f, indent=2)
            
            # Copy original file
            shutil.copy2(file_path, module_output_dir / file_path.name)

            if not args.dry_run:
                # Upload to S3
                s3_base = f"{args.s3_prefix}/{module_code}" if args.s3_prefix else module_code
                
                # Upload everything in module_output_dir
                for path in module_output_dir.rglob('*'):
                    if path.is_file():
                        rel_path = path.relative_to(module_output_dir)
                        s3_key = f"{s3_base}/{rel_path}".replace('\\', '/')
                        print(f"Uploading {path} to s3://{args.bucket}/{s3_key}")
                        s3.upload_file(str(path), args.bucket, s3_key)
                
                # Cleanup temp
                shutil.rmtree(module_output_dir)
            
            materials.append(material_data)
            
        except Exception as e:
            print(f"Error processing {file_path.name}: {e}")
            import traceback
            traceback.print_exc()

    # If dry run, save index
    if args.dry_run:
        with open(output_base / 'index.json', 'w') as f:
            json.dump(materials, f, indent=2)
        print(f"Dry run complete. Output in {args.output_dir}")

if __name__ == '__main__':
    main()
